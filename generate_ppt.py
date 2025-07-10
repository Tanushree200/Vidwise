from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from io import BytesIO
import google.generativeai as genai
import os
import json
from dotenv import load_dotenv


load_dotenv()
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))


def generate_slide_data(transcript_text: str):
    prompt = f"""
You are a presentation expert.

Using the full transcript below, generate a detailed slide deck structure.
Break it into 5–7 slides. Each slide should have:
- A relevant title
- 3–5 bullet points (clear, informative)

Return ONLY JSON like:
[
  {{
    "title": "Slide Title",
    "bullets": ["Bullet 1", "Bullet 2"]
  }}
]
TRANSCRIPT:
\"\"\"{transcript_text[:12000]}\"\"\"
"""
    model = genai.GenerativeModel("gemini-2.0-flash")
    response = model.generate_content(prompt)
    raw = response.text.strip()

    import re
    try:
        cleaned = re.sub(r"^```(?:json)?|```$", "", raw, flags=re.MULTILINE).strip()
        slides = json.loads(cleaned)
    except Exception as e:
        slides = [{
            "title": "Error Slide",
            "bullets": [f"⚠️ Failed to parse Gemini response: {str(e)}"]
        }]
    return slides



def create_ppt_from_summary(summary_text: str, save_path: str = None) -> BytesIO:
    prs = Presentation()

    def apply_background(slide, color=RGBColor(245, 245, 255)):  
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def style_title(title_shape):
        title_frame = title_shape.text_frame
        p = title_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.name = 'Segoe UI'
        p.font.color.rgb = RGBColor(40, 40, 40)

    def style_bullets(content_shape, bullets):
        tf = content_shape.text_frame
        tf.clear()
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.level = 0
            p.font.size = Pt(20)
            p.font.name = 'Segoe UI'
            p.font.color.rgb = RGBColor(60, 60, 60)

    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_background(slide)
    slide.shapes.title.text = "VidWise"
    style_title(slide.shapes.title)
    subtitle = slide.placeholders[1]
    subtitle.text = "Auto-generated Presentation"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(90, 90, 90)

    
    slides_data = generate_slide_data(summary_text)

    for data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        apply_background(slide)
        slide.shapes.title.text = data.get("title", "Untitled Slide")
        style_title(slide.shapes.title)
        style_bullets(slide.placeholders[1], data.get("bullets", []))

    
    recap = prs.slides.add_slide(prs.slide_layouts[1])
    apply_background(recap)
    recap.shapes.title.text = "🔁 Recap of Topics"
    style_title(recap.shapes.title)
    style_bullets(recap.placeholders[1], [s["title"] for s in slides_data])

    
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    if save_path:
        os.makedirs(os.path.dirname(save_path), exist_ok=True)
        with open(save_path, "wb") as f:
            f.write(ppt_stream.getvalue())

    return ppt_stream
